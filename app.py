from flask import Flask, render_template, request, redirect, url_for, flash, session, jsonify, Response
from flask_sqlalchemy import SQLAlchemy
from flask_socketio import SocketIO, join_room, emit
from werkzeug.security import generate_password_hash, check_password_hash
from groq import Groq
import json, os, datetime, uuid, requests

app = Flask(__name__)
app.config['SECRET_KEY'] = 'zyntra-secret-2025'
app.config['SQLALCHEMY_DATABASE_URI'] = 'sqlite:///zyntra.db'
app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False
db = SQLAlchemy(app)
socketio = SocketIO(app, cors_allowed_origins="*", async_mode='threading')

class User(db.Model):
    id       = db.Column(db.Integer, primary_key=True)
    username = db.Column(db.String(80), unique=True, nullable=False)
    password = db.Column(db.String(200), nullable=False)
    memory   = db.Column(db.Text, default='{}')
    chats    = db.relationship('Chat', backref='user', lazy=True, cascade='all,delete')

class Chat(db.Model):
    id          = db.Column(db.Integer, primary_key=True)
    name        = db.Column(db.String(100), default='Nouveau chat')
    messages    = db.Column(db.Text, default='[]')
    user_id     = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    created     = db.Column(db.DateTime, default=datetime.datetime.utcnow)
    is_pinned   = db.Column(db.Boolean, default=False)
    is_archived = db.Column(db.Boolean, default=False)

class Group(db.Model):
    id          = db.Column(db.Integer, primary_key=True)
    name        = db.Column(db.String(100), nullable=False)
    invite_code = db.Column(db.String(32), unique=True, nullable=False)
    created_by  = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    created     = db.Column(db.DateTime, default=datetime.datetime.utcnow)
    messages    = db.Column(db.Text, default='[]')

class GroupMember(db.Model):
    id       = db.Column(db.Integer, primary_key=True)
    group_id = db.Column(db.Integer, db.ForeignKey('group.id'), nullable=False)
    user_id  = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    joined   = db.Column(db.DateTime, default=datetime.datetime.utcnow)

def get_client():
    try:
        from config import API_KEY
    except:
        API_KEY = ''
    return Groq(api_key=API_KEY)

def current_user():
    if 'user_id' not in session:
        return None
    return User.query.get(session['user_id'])

def extract_memory(text, memory):
    t = text.lower()
    if "je m'appelle" in t:
        try:
            name = text.split("je m'appelle")[-1].strip().split()[0]
            if 2 <= len(name) <= 20:
                memory['prenom'] = name.capitalize()
        except:
            pass
    return memory

def get_image(query):
    try:
        from config import PEXELS_KEY
        headers = {'Authorization': PEXELS_KEY}
        url = f"https://api.pexels.com/v1/search?query={query}&per_page=5&orientation=landscape"
        r = requests.get(url, headers=headers, timeout=10)
        data = r.json()
        photos = data.get('photos', [])
        if not photos:
            return None
        img_url = photos[0]['src']['large']
        img_data = requests.get(img_url, timeout=15).content
        static_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'static')
        os.makedirs(static_dir, exist_ok=True)
        path = os.path.join(static_dir, f'img_{uuid.uuid4().hex[:8]}.jpg')
        with open(path, 'wb') as f:
            f.write(img_data)
        return path
    except Exception as e:
        print(f"Erreur image: {e}")
        return None

def make_pptx(slides_data, filename):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    DARK   = RGBColor(0x0d, 0x0d, 0x0d)
    ORANGE = RGBColor(0xCC, 0x78, 0x5C)
    WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
    GREY   = RGBColor(0xAA, 0xAA, 0xAA)
    LIGHT  = RGBColor(0xE0, 0xE0, 0xE0)

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    for i, slide_data in enumerate(slides_data):
        layout = prs.slide_layouts[6]
        slide  = prs.slides.add_slide(layout)
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = DARK

        if i == 0:
            img_path = get_image(slide_data.get('title', 'technology'))
            if img_path and os.path.exists(img_path):
                slide.shapes.add_picture(img_path, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
            overlay = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
            overlay.fill.solid()
            overlay.fill.fore_color.rgb = RGBColor(0x0d, 0x0d, 0x0d)
            overlay.line.fill.background()
            sp_tree = slide.shapes._spTree
            sp_tree.remove(overlay._element)
            sp_tree.insert(3, overlay._element)
            box = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.3), Inches(2))
            tf  = box.text_frame
            tf.word_wrap = True
            p   = tf.paragraphs[0]
            p.text = slide_data.get('title', '')
            p.alignment = PP_ALIGN.CENTER
            if p.runs:
                p.runs[0].font.size  = Pt(52)
                p.runs[0].font.bold  = True
                p.runs[0].font.color.rgb = ORANGE
            sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.3), Inches(1))
            sf = sub_box.text_frame
            sp2 = sf.paragraphs[0]
            sp2.text = slide_data.get('subtitle', 'Powered by Zyntra')
            sp2.alignment = PP_ALIGN.CENTER
            if sp2.runs:
                sp2.runs[0].font.size = Pt(22)
                sp2.runs[0].font.color.rgb = GREY
            line = slide.shapes.add_shape(1, Inches(4), Inches(4.2), Inches(5.3), Inches(0.05))
            line.fill.solid()
            line.fill.fore_color.rgb = ORANGE
            line.line.fill.background()
        else:
            header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.3))
            header.fill.solid()
            header.fill.fore_color.rgb = ORANGE
            header.line.fill.background()
            num = slide.shapes.add_textbox(Inches(0.2), Inches(0.1), Inches(0.8), Inches(1.1))
            nf  = num.text_frame
            np2 = nf.paragraphs[0]
            np2.text = str(i)
            if np2.runs:
                np2.runs[0].font.size = Pt(36)
                np2.runs[0].font.bold = True
                np2.runs[0].font.color.rgb = WHITE
            title_box = slide.shapes.add_textbox(Inches(1), Inches(0.15), Inches(11.5), Inches(1))
            ttf = title_box.text_frame
            tp  = ttf.paragraphs[0]
            tp.text = slide_data.get('title', '')
            if tp.runs:
                tp.runs[0].font.size = Pt(30)
                tp.runs[0].font.bold = True
                tp.runs[0].font.color.rgb = WHITE
            img_path = get_image(slide_data.get('title', 'technology'))
            if img_path and os.path.exists(img_path):
                slide.shapes.add_picture(img_path, Inches(8.8), Inches(1.4), Inches(4.2), Inches(5.5))
            content_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(8.0), Inches(5.7))
            ctf = content_box.text_frame
            ctf.word_wrap = True
            ctf.text = ''
            for j, point in enumerate(slide_data.get('points', [])):
                p = ctf.add_paragraph() if j > 0 else ctf.paragraphs[0]
                p.text = '▸   ' + point
                p.space_before = Pt(10)
                if p.runs:
                    p.runs[0].font.size = Pt(18)
                    p.runs[0].font.color.rgb = LIGHT
            footer = slide.shapes.add_shape(1, Inches(0), Inches(7.2), Inches(13.33), Inches(0.3))
            footer.fill.solid()
            footer.fill.fore_color.rgb = RGBColor(0x1a, 0x1a, 0x1a)
            footer.line.fill.background()
            ft = slide.shapes.add_textbox(Inches(0.2), Inches(7.2), Inches(13), Inches(0.3))
            ff = ft.text_frame
            fp = ff.paragraphs[0]
            fp.text = 'Zyntra AI'
            fp.alignment = PP_ALIGN.RIGHT
            if fp.runs:
                fp.runs[0].font.size = Pt(10)
                fp.runs[0].font.color.rgb = GREY

    static_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'static')
    filepath = os.path.join(static_dir, filename)
    prs.save(filepath)

def detect_expo_request(text):
    keywords = ['powerpoint', 'présentation', 'presentation', 'ppt', 'expo', 'diaporama', 'slides', 'exposé', 'expose']
    return any(k in text.lower() for k in keywords)

def generate_chat_title(text):
    try:
        client = get_client()
        resp = client.chat.completions.create(
            model='llama-3.3-70b-versatile',
            messages=[
                {'role': 'system', 'content': "Génère un titre court de 3 à 5 mots maximum qui résume le sujet de ce message. Réponds UNIQUEMENT avec le titre, sans guillemets, sans ponctuation finale, rien d'autre."},
                {'role': 'user', 'content': text}
            ],
            max_tokens=20
        )
        title = resp.choices[0].message.content.strip().strip('"').strip("'")
        return title[:40] if title else ' '.join(text.split()[:5])[:35]
    except:
        return ' '.join(text.split()[:5])[:35]

@app.route('/')
def index():
    if 'user_id' not in session:
        return redirect(url_for('login'))
    return redirect(url_for('chat'))

@app.route('/login', methods=['GET','POST'])
def login():
    if request.method == 'POST':
        username = request.form.get('username','').strip()
        password = request.form.get('password','')
        user = User.query.filter_by(username=username).first()
        if user and check_password_hash(user.password, password):
            session['user_id'] = user.id
            return redirect(url_for('chat'))
        flash('Nom ou mot de passe incorrect.')
    return render_template('login.html')

@app.route('/register', methods=['GET','POST'])
def register():
    if request.method == 'POST':
        username = request.form.get('username','').strip()
        password = request.form.get('password','')
        if not username or not password:
            flash('Remplis tous les champs.')
        elif User.query.filter_by(username=username).first():
            flash('Ce nom est deja pris.')
        else:
            user = User(username=username, password=generate_password_hash(password))
            db.session.add(user)
            db.session.commit()
            session['user_id'] = user.id
            return redirect(url_for('chat'))
    return render_template('register.html')

@app.route('/logout')
def logout():
    session.clear()
    return redirect(url_for('login'))

@app.route('/chat')
def chat():
    user = current_user()
    if not user:
        return redirect(url_for('login'))
    chats = Chat.query.filter_by(user_id=user.id, is_archived=False).order_by(Chat.is_pinned.desc(), Chat.created.desc()).all()
    if not chats:
        c = Chat(name='Nouveau chat', user_id=user.id)
        db.session.add(c)
        db.session.commit()
        chats = [c]
    active_id = request.args.get('chat_id', chats[0].id)
    active = Chat.query.get(active_id) or chats[0]
    messages = json.loads(active.messages)
    memory = json.loads(user.memory)
    groups = db.session.query(Group).join(GroupMember, Group.id == GroupMember.group_id).filter(GroupMember.user_id == user.id).all()
    return render_template('chat.html', user=user, chats=chats, active=active, messages=messages, memory=memory, groups=groups)

@app.route('/group/<int:group_id>')
def group_chat(group_id):
    user = current_user()
    if not user:
        return redirect(url_for('login'))
    group = Group.query.get_or_404(group_id)
    member = GroupMember.query.filter_by(group_id=group_id, user_id=user.id).first()
    if not member:
        return redirect(url_for('chat'))
    messages = json.loads(group.messages)
    members = db.session.query(User).join(GroupMember, User.id == GroupMember.user_id).filter(GroupMember.group_id == group_id).all()
    chats = Chat.query.filter_by(user_id=user.id, is_archived=False).order_by(Chat.is_pinned.desc(), Chat.created.desc()).all()
    groups = db.session.query(Group).join(GroupMember, Group.id == GroupMember.group_id).filter(GroupMember.user_id == user.id).all()
    return render_template('group.html', user=user, group=group, messages=messages, members=members, chats=chats, groups=groups)

@app.route('/join/<invite_code>')
def join_group(invite_code):
    user = current_user()
    if not user:
        session['next'] = f'/join/{invite_code}'
        return redirect(url_for('login'))
    group = Group.query.filter_by(invite_code=invite_code).first()
    if not group:
        flash('Lien invalide.')
        return redirect(url_for('chat'))
    existing = GroupMember.query.filter_by(group_id=group.id, user_id=user.id).first()
    if not existing:
        member = GroupMember(group_id=group.id, user_id=user.id)
        db.session.add(member)
        db.session.commit()
    return redirect(url_for('group_chat', group_id=group.id))

@app.route('/api/new_chat', methods=['POST'])
def new_chat():
    user = current_user()
    if not user: return jsonify({'error':'non connecte'}), 401
    c = Chat(name='Nouveau chat', user_id=user.id)
    db.session.add(c)
    db.session.commit()
    return jsonify({'id': c.id, 'name': c.name})

@app.route('/api/delete_chat', methods=['POST'])
def delete_chat():
    user = current_user()
    if not user: return jsonify({'error':'non connecte'}), 401
    data = request.json
    c = Chat.query.get(data.get('chat_id'))
    if c and c.user_id == user.id:
        db.session.delete(c)
        db.session.commit()
    return jsonify({'ok': True})

@app.route('/api/rename_chat', methods=['POST'])
def rename_chat():
    user = current_user()
    if not user: return jsonify({'error':'non connecte'}), 401
    data = request.json
    c = Chat.query.get(data.get('chat_id'))
    if c and c.user_id == user.id:
        c.name = data.get('name', c.name)
        db.session.commit()
    return jsonify({'ok': True})

@app.route('/api/pin_chat', methods=['POST'])
def pin_chat():
    user = current_user()
    if not user: return jsonify({'error':'non connecte'}), 401
    data = request.json
    c = Chat.query.get(data.get('chat_id'))
    if c and c.user_id == user.id:
        c.is_pinned = not c.is_pinned
        db.session.commit()
    return jsonify({'ok': True, 'pinned': c.is_pinned})

@app.route('/api/archive_chat', methods=['POST'])
def archive_chat():
    user = current_user()
    if not user: return jsonify({'error':'non connecte'}), 401
    data = request.json
    c = Chat.query.get(data.get('chat_id'))
    if c and c.user_id == user.id:
        c.is_archived = True
        db.session.commit()
    return jsonify({'ok': True})

@app.route('/api/share_chat', methods=['POST'])
def share_chat():
    user = current_user()
    if not user: return jsonify({'error':'non connecte'}), 401
    data = request.json
    c = Chat.query.get(data.get('chat_id'))
    if not c or c.user_id != user.id:
        return jsonify({'error':'introuvable'}), 404
    share_token = uuid.uuid4().hex[:12]
    return jsonify({'ok': True, 'link': f'/share/{share_token}', 'full_link': f'http://127.0.0.1:5000/share/{share_token}'})

@app.route('/api/create_group', methods=['POST'])
def create_group():
    user = current_user()
    if not user: return jsonify({'error':'non connecte'}), 401
    data = request.json
    name = data.get('name', 'Groupe Zyntra')
    invite_code = uuid.uuid4().hex[:10]
    group = Group(name=name, invite_code=invite_code, created_by=user.id)
    db.session.add(group)
    db.session.commit()
    member = GroupMember(group_id=group.id, user_id=user.id)
    db.session.add(member)
    db.session.commit()
    return jsonify({'id': group.id, 'name': group.name, 'invite_code': invite_code, 'invite_link': f'http://127.0.0.1:5000/join/{invite_code}'})

@app.route('/api/delete_group', methods=['POST'])
def delete_group():
    user = current_user()
    if not user: return jsonify({'error':'non connecte'}), 401
    data = request.json
    g = Group.query.get(data.get('group_id'))
    if g and g.created_by == user.id:
        GroupMember.query.filter_by(group_id=g.id).delete()
        db.session.delete(g)
        db.session.commit()
    return jsonify({'ok': True})

@socketio.on('join_group')
def on_join(data):
    join_room(f"group_{data['group_id']}")

@socketio.on('group_message')
def on_group_message(data):
    user = current_user()
    if not user: return
    group_id = data.get('group_id')
    text = data.get('message', '').strip()
    if not text: return
    group = Group.query.get(group_id)
    if not group: return
    member = GroupMember.query.filter_by(group_id=group_id, user_id=user.id).first()
    if not member: return
    now = datetime.datetime.now().strftime('%H:%M')
    messages = json.loads(group.messages)
    messages.append({'role': 'user', 'sender': user.username, 'content': text, 'time': now})
    try:
        client = get_client()
        sys_msg = f"You are Zyntra, an AI assistant in a group chat. The group is called '{group.name}'. Always respond in the same language as the user."
        history = [{'role': 'user' if m['role'] == 'user' else 'assistant', 'content': m['content']} for m in messages[-6:]]
        response = client.chat.completions.create(model='llama-3.3-70b-versatile', messages=[{'role': 'system', 'content': sys_msg}] + history, max_tokens=1024)
        reply = response.choices[0].message.content.strip()
    except Exception as e:
        reply = f"Erreur : {str(e)}"
    messages.append({'role': 'assistant', 'sender': 'Zyntra', 'content': reply, 'time': now})
    group.messages = json.dumps(messages, ensure_ascii=False)
    db.session.commit()
    room = f"group_{group_id}"
    emit('new_message', {'sender': user.username, 'content': text, 'time': now, 'role': 'user'}, room=room)
    emit('new_message', {'sender': 'Zyntra', 'content': reply, 'time': now, 'role': 'assistant'}, room=room)

@app.route('/api/send', methods=['POST'])
def send():
    user = current_user()
    if not user: return jsonify({'error': 'non connecte'}), 401
    data    = request.json
    chat_id = data.get('chat_id')
    text    = data.get('message', '').strip()
    mode    = data.get('mode', 'Normal')
    if not text: return jsonify({'error': 'message vide'}), 400
    c = Chat.query.get(chat_id)
    if not c or c.user_id != user.id:
        return jsonify({'error': 'chat introuvable'}), 404

    messages = json.loads(c.messages)
    memory   = json.loads(user.memory)

    if len(messages) == 0:
        c.name = generate_chat_title(text)
        db.session.commit()

    memory = extract_memory(text, memory)
    user.memory = json.dumps(memory, ensure_ascii=False)
    now = datetime.datetime.now().strftime('%H:%M')
    messages.append({'role': 'user', 'content': text, 'time': now})

    is_expo = mode == 'Expo' or detect_expo_request(text)

    if is_expo:
        sys_msg = f"""You are a PowerPoint content generator. The user wants a presentation about: {text}
Respond in the same language as the user. Use EXACTLY this format:

TITLE: Main title here
SUBTITLE: Subtitle here

SLIDE: Slide 2 title
- Full detailed sentence of minimum 10 words.
- Full detailed sentence of minimum 10 words.
- Full detailed sentence of minimum 10 words.
- Full detailed sentence of minimum 10 words.

SLIDE: Slide 3 title
- Full detailed sentence.
- Full detailed sentence.
- Full detailed sentence.

SLIDE: Slide 4 title
- Full detailed sentence.
- Full detailed sentence.
- Full detailed sentence.

SLIDE: Slide 5 title
- Full detailed sentence.
- Full detailed sentence.
- Full detailed sentence.

SLIDE: Conclusion
- Key takeaway 1.
- Key takeaway 2.
- Key takeaway 3.
- Final thought."""
    elif mode == 'Expert Code':
        sys_msg = "You are Zyntra, a programming expert. Give clean, commented code with clear explanations. Always respond in the same language as the user."
    else:
        sys_msg = "You are Zyntra, an intelligent AI assistant. Always respond in the same language as the user."

    if memory.get('prenom'):
        sys_msg += f" Tu parles a {memory['prenom']}."

    history = [{'role': m['role'], 'content': m['content'][:400]} for m in messages[-6:]]

    if is_expo:
        try:
            client = get_client()
            response = client.chat.completions.create(model='llama-3.3-70b-versatile', messages=[{'role': 'system', 'content': sys_msg}] + history, max_tokens=2048)
            raw = response.choices[0].message.content.strip()
            slides_data = []
            lines = raw.strip().split('\n')
            current_slide = None
            title = ''
            subtitle = ''
            for line in lines:
                line = line.strip()
                if line.startswith('TITLE:'):
                    title = line.replace('TITLE:', '').strip()
                elif line.startswith('SUBTITLE:'):
                    subtitle = line.replace('SUBTITLE:', '').strip()
                    slides_data.append({'title': title, 'subtitle': subtitle})
                elif line.startswith('SLIDE:'):
                    if current_slide:
                        slides_data.append(current_slide)
                    current_slide = {'title': line.replace('SLIDE:', '').strip(), 'points': []}
                elif line.startswith('- ') and current_slide:
                    current_slide['points'].append(line[2:])
            if current_slide:
                slides_data.append(current_slide)
            filename = f'zyntra_expo_{uuid.uuid4().hex[:8]}.pptx'
            make_pptx(slides_data, filename)
            reply = f'Votre presentation PowerPoint est prete !\n\nTelecharger le PowerPoint : /static/{filename}'
        except Exception as e:
            reply = f'Erreur creation PowerPoint : {str(e)}'
        messages.append({'role': 'assistant', 'content': reply, 'time': now})
        c.messages = json.dumps(messages, ensure_ascii=False)
        db.session.commit()
        return jsonify({'reply': reply, 'chat_name': c.name, 'time': now})

    # ✅ STREAMING lettre par lettre
    def generate():
        full_reply = ''
        try:
            client = get_client()
            stream = client.chat.completions.create(
                model='llama-3.3-70b-versatile',
                messages=[{'role': 'system', 'content': sys_msg}] + history,
                max_tokens=5000,
                stream=True
            )
            for chunk in stream:
                token = chunk.choices[0].delta.content or ''
                if token:
                    full_reply += token
                    yield f"data: {json.dumps({'token': token})}\n\n"
        except Exception as e:
            err = f'Erreur: {str(e)}'
            full_reply = err
            yield f"data: {json.dumps({'token': err})}\n\n"

        messages.append({'role': 'assistant', 'content': full_reply, 'time': now})
        c.messages = json.dumps(messages, ensure_ascii=False)
        db.session.commit()
        yield f"data: {json.dumps({'done': True, 'chat_name': c.name, 'time': now})}\n\n"

    return Response(generate(), mimetype='text/event-stream')

if __name__ == '__main__':
    with app.app_context():
        db.create_all()
    socketio.run(app, debug=True, port=5000, allow_unsafe_werkzeug=True)